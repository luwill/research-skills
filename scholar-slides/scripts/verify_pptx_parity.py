#!/usr/bin/env python3
"""HTML/PDF <-> PPTX content-parity regression.

The PPTX is the editable deliverable; the promise is "minimal manual edits". That breaks the
moment the exporter silently drops or mangles content — a lost annotation, a mangled table cell,
missing speaker notes — because the user then has to repair it by hand. There is no LibreOffice
here to pixel-diff, and a pixel diff would be flaky anyway; what actually matters is that every
load-bearing element in the deck spec survives into the PPTX *natively*.

This checks exactly that, per slide:
  - slide count matches
  - every expected text string (titles, bullets, annotations, table cells, captions, questions,
    references, section headers, cover metadata) appears in the slide's native text
  - speaker notes survive as native notes
  - a slide with a figure has a picture; a results-table slide has a native table

Figures and equations are images by design (the honest, flagged degradation) — parity requires the
image to be PRESENT, not that it carry text. The comparison core (`check_parity`) is pure and
unit-tested; `extract_from_pptx` is the thin python-pptx adapter.
"""
from __future__ import annotations

import json
import re
import sys

_MATH = re.compile(r"\$([^$]+)\$")
_EMPH = re.compile(r"==(?:[kpw]\|)?(.+?)==")


def strip_math(s: str) -> str:
    """Mirror pptx_layout.stripInlineMath: drop $...$ math and ==...== emphasis, keep inner text."""
    return _EMPH.sub(r"\1", _MATH.sub(r"\1", str(s or "")))


def norm(s: str) -> str:
    """Whitespace-collapsed, math/emphasis-stripped, lowercased — for a presence check."""
    return " ".join(strip_math(s).split()).lower()


def contains(haystack_norm: str, want_norm: str) -> bool:
    """Boundary-aware presence check on already-norm()ed strings. Plain substring matching let
    an unrelated "59" satisfy dropped cells "5" and "9"; here the match may not sit inside an
    ASCII-alphanumeric run (CJK has no word delimiters, so CJK containment is unaffected)."""
    if not want_norm:
        return True
    return re.search(
        r"(?<![0-9a-z])" + re.escape(want_norm) + r"(?![0-9a-z])", haystack_norm
    ) is not None


def _cells(table) -> list[str]:
    out = []
    for row in table.get("rows", []) or []:
        for cell in row:
            out.append(cell.get("v") if isinstance(cell, dict) else cell)
    return out


def expected_texts(s) -> list[tuple[str, str]]:
    """(field, text) strings the exporter emits as NATIVE text for this slide. Mirrors
    export_pptx.buildSlide — keep in sync when the exporter's text contract changes."""
    L = s.get("layout")
    out: list[tuple[str, str]] = []

    def add(field, text):
        if text is not None and str(text).strip():
            out.append((field, str(text)))

    if L == "paper-title":
        add("title", s.get("title"))
        a = s.get("authors")
        add("authors", ", ".join(a) if isinstance(a, list) else a)
        add("affiliation", s.get("affiliation"))
        add("venue", s.get("venue"))
        add("presenter", s.get("presenter"))
        return out
    if L == "section":
        add("num", s.get("num"))
        add("title", s.get("title"))
        return out

    # every other layout goes through titleBar
    add("eyebrow", s.get("eyebrow"))
    add("title", s.get("action_title") or s.get("title"))

    if L == "bullets":
        for i, p in enumerate(s.get("points", []) or []):
            add(f"points[{i}]", p)
    elif L == "outline-agenda":
        for i, it in enumerate(s.get("items", []) or []):
            add(f"items[{i}]", it)
    elif L == "assertion-evidence":
        if s.get("figure") and s["figure"].get("caption"):
            add("figure.caption", s["figure"]["caption"])
        add("annotation", s.get("annotation"))
    elif L == "equation":
        add("note", s.get("note"))
    elif L == "results-table":
        t = s.get("table", {}) or {}
        add("table.caption", t.get("caption"))
        for i, c in enumerate(_cells(t)):
            add(f"table.cell[{i}]", c)
        add("table.footnote", t.get("footnote"))
    elif L == "two-column":
        for i, p in enumerate(s.get("points", []) or []):
            add(f"points[{i}]", p)
        add("text", s.get("text"))
        for i, p in enumerate(s.get("points2", []) or []):
            add(f"points2[{i}]", p)
    elif L == "critique-concerns":
        for i, p in enumerate(s.get("points", []) or []):
            add(f"points[{i}].head", p.get("head"))
            add(f"points[{i}].body", p.get("body"))
    elif L == "discussion-questions":
        for i, q in enumerate(s.get("questions", []) or []):
            add(f"questions[{i}]", q)
    elif L == "references":
        for i, e in enumerate(s.get("entries", []) or []):
            add(f"entries[{i}]", e)
    return out


def expects_figure(s) -> bool:
    return bool(s.get("figure") and s["figure"].get("src"))


def expects_table(s) -> bool:
    return s.get("layout") == "results-table" and bool(s.get("table"))


def check_parity(deck, extracted) -> list[dict]:
    """Pure comparison. `extracted` = per-slide dicts {text, notes, has_picture, has_table}
    (from extract_from_pptx). Returns a list of parity findings (empty = faithful)."""
    findings = []
    slides = deck.get("slides", []) or []
    if len(slides) != len(extracted):
        findings.append({"slide": 0, "kind": "slide-count",
                         "detail": f"spec has {len(slides)} slides, pptx has {len(extracted)}"})
    for i, (s, ex) in enumerate(zip(slides, extracted)):
        text = norm(ex.get("text", ""))
        # When the adapter supplies per-cell texts, table cells are checked cell-against-cell and
        # CONSUMED (multiset): one native "5" cannot vouch for two expected "5"s, and a stray
        # number elsewhere on the slide cannot vouch for a dropped cell at all.
        cells = ex.get("cells")
        remaining = [norm(c) for c in cells] if cells is not None else None
        for field, want in expected_texts(s):
            w = norm(want)
            if field.startswith("table.cell[") and remaining is not None:
                hit = next((j for j, c in enumerate(remaining) if contains(c, w)), None)
                if hit is None:
                    findings.append({"slide": i + 1, "kind": "missing-text", "field": field,
                                     "detail": f'"{" ".join(str(want).split())[:70]}" not native in pptx'})
                else:
                    remaining.pop(hit)
                continue
            if not contains(text, w):
                findings.append({"slide": i + 1, "kind": "missing-text", "field": field,
                                 "detail": f'"{" ".join(str(want).split())[:70]}" not native in pptx'})
        notes = s.get("speaker_notes")
        if notes and not contains(norm(ex.get("notes", "")), norm(notes)):
            findings.append({"slide": i + 1, "kind": "missing-notes",
                             "detail": "speaker notes not native in pptx"})
        if expects_figure(s) and not ex.get("has_picture"):
            findings.append({"slide": i + 1, "kind": "missing-figure",
                             "detail": f'figure {s["figure"]["src"]} has no picture shape'})
        if expects_table(s) and not ex.get("has_table"):
            findings.append({"slide": i + 1, "kind": "missing-table",
                             "detail": "results-table slide has no native table"})
    return findings


def extract_from_pptx(pptx_path) -> list[dict]:
    """Adapter: read a .pptx into the per-slide shape python check_parity expects."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    out = []
    for slide in prs.slides:
        parts, cells, has_pic, has_tbl = [], [], False, False
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if sh.has_table:
                has_tbl = True
                for row in sh.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
                        cells.append(cell.text)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_pic = True
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        out.append({"text": " \n ".join(parts), "notes": notes, "cells": cells,
                    "has_picture": has_pic, "has_table": has_tbl})
    return out


def main(argv):
    import argparse

    ap = argparse.ArgumentParser(description="Verify PPTX preserves the deck spec's content natively.")
    ap.add_argument("deck_json")
    ap.add_argument("pptx")
    args = ap.parse_args(argv)
    deck = json.load(open(args.deck_json, encoding="utf-8"))
    findings = check_parity(deck, extract_from_pptx(args.pptx))
    if not findings:
        print(f"PPTX parity OK — {len(deck.get('slides', []))} slides, all native content preserved.")
        return 0
    print(f"PPTX parity: {len(findings)} issue(s)")
    for f in findings:
        print(f"  [slide {f['slide']}] {f['kind']}: {f.get('field', '')} {f['detail']}")
    return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
